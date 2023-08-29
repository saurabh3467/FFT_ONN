# -*- coding: utf-8 -*-
"""
Created on Sun Jun 18 16:13:28 2023

@author: velle
"""

import os
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib.pyplot as plt

# Ask the user to input the folder path
folder_path = input("Enter the path to the folder: ")

# Ask the user to specify the column indices
time_column = 0
current_column = 1

# Ask the user to enter the interval for the zoomed-in plot
interval = 10

# Get a list of files in the folder
file_list = [file for file in os.listdir(folder_path) if file.lower().endswith(".txt")]

# Create a new PowerPoint presentation
presentation = Presentation()

# Loop through each file in the folder
for file_name in file_list:
    # Construct the file path
    file_path = os.path.join(folder_path, file_name)

    try:
        # Read the data from the text file, skipping the first row
        data = np.loadtxt(file_path, skiprows=1)

        # Check if the loaded data has at least two rows
        if data.shape[0] >= 2:
            # Extract the time and current columns
            time = data[:, time_column]
            current = data[:, current_column]

            # Plot time vs current
            fig, axs = plt.subplots(1, 3, figsize=(10, 4))

            # Plot 1: Time vs Current
            axs[0].plot(time, current, linewidth=0.5)
            axs[0].set_xlabel("Time")
            axs[0].set_ylabel("Current")
            axs[0].set_title("Time vs Current")

            # Zoom in on the plot
            start_time = time[0] + 30
            end_time = start_time + interval

            # Filter the current array based on the time range
            filtered_current = current[(time >= start_time) & (time <= end_time)]

            # Find the minimum and maximum y-values within the time range
            min_y = np.min(filtered_current)
            max_y = np.max(filtered_current)

            # Plot 2: Zoomed-in Time vs Current
            axs[1].plot(time, current, linewidth=0.5)
            axs[1].set_xlabel("Time")
            axs[1].set_ylabel("Current")
            axs[1].set_title("Zoomed-in Time vs Current")
            axs[1].set_xlim(start_time, end_time)
            axs[1].set_ylim(min_y-0.0005, max_y+0.0005) # set zoomed in values, adjust if the zoomin goes weirdly

            # Perform Fast Fourier Transform (FFT)
            window = np.hanning(len(current))
            current_windowed = current * window
            fft_current = np.fft.fft(current)
# time = [] , length(time), time [-1] - time [0] / length
            frequencies = np.fft.fftfreq(len(current), d=(time[25] - time[24])) # Change the sampling rate calculation to that of origins methodology it should work.

            # Filter the FFT result and frequencies for positive frequencies
            positive_mask = frequencies > 1
            positive_frequencies = frequencies[positive_mask]
            positive_fft_result = fft_current[positive_mask]

            # Plot 3: Fast Fourier Transform (FFT)
            axs[2].plot(positive_frequencies, np.abs(positive_fft_result), linewidth=1, color='red')
            axs[2].set_xlabel("Frequency")
            axs[2].set_ylabel("Amplitude")
            axs[2].set_title("Fast Fourier Transform (FFT)")
            axs[2].set_xlim(-1, 15)

            # Calculate the maximum amplitude in the positive frequency range
            max_amplitude = np.max(np.abs(positive_fft_result))

            # Set the y-limit for the FFT plot
            # axs[2].set_ylim(0, 1.1 * max_amplitude)

            # Adjust spacing between subplots
            plt.subplots_adjust(wspace=0.4)

            # Save the figure to an image file
            plots_file_path = os.path.join(folder_path, "plots.png")
            plt.savefig(plots_file_path, dpi=300)

            # Create a new PowerPoint slide
            slide_layout = presentation.slide_layouts[5]
            slide = presentation.slides.add_slide(slide_layout)

            # Add the file name to the slide
            title = slide.shapes.title
            title.text = file_name

            left = Inches(1)
            top = Inches(1.5)
            slide.shapes.add_picture(plots_file_path, left, top)

            # Add the highest peak in the higher than zero range to the slide
            peak_index = np.argmax(np.abs(positive_fft_result))
            peak_frequency = positive_frequencies[peak_index]
            text_box = slide.shapes.add_textbox(left + Inches(5), top, width=Inches(2), height=Inches(1))
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = "Highest Peak (f > 0 Hz):"
            p = text_frame.add_paragraph()
            p.text = "{:.2f} Hz".format(peak_frequency)

            # Find the highest peak in the area of 15 Hz and above
            # To get the frequency from the cracking oscillations
            higher_than_15_mask = positive_frequencies >= 15
            higher_than_15_fft_result = positive_fft_result[higher_than_15_mask]
            higher_than_15_frequencies = positive_frequencies[higher_than_15_mask]
            highest_peak_index = np.argmax(np.abs(higher_than_15_fft_result))
            highest_peak_frequency = higher_than_15_frequencies[highest_peak_index]

            # Add the highest peak in the area of 15 Hz and above to the slide
            p = text_frame.add_paragraph()
            p.text = "Highest Peak (f >= 15 Hz):"
            p = text_frame.add_paragraph()
            p.text = "{:.2f} Hz".format(highest_peak_frequency)

            # Save the PowerPoint presentation
            presentation_file_name = os.path.splitext(file_name)[0] + ".pptx"
            presentation_file_path = os.path.join(folder_path, presentation_file_name)
            presentation.save(presentation_file_path)

            # Show a success message
            print(f"The PowerPoint presentation has been created: {presentation_file_path}")
        else:
            print(f"The file {file_name} does not have enough data points.")
    except ValueError:
        print(f"Error reading data from file {file_name}. Skipping file.")
